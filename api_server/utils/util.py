import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import io
from sentence_transformers import SentenceTransformer
import psycopg2
import dotenv
import os
import docx
from pptx import Presentation
import openpyxl
import xlrd
import re

dotenv.load_dotenv()

model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')

def extract_pdf_text(file_storage) -> str:
    """
    Extract text from a PDF using PyMuPDF (fitz).
    Accepts a Flask FileStorage object (e.g. from request.files).
    """
    # Load the PDF from a byte stream (not from disk)
    pdf_bytes = file_storage.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")

    text = ""
    for page in doc:
        text += page.get_text() + "\n"

    doc.close()
    return text.strip()


def extract_image_text(file_storage) -> str:
    """
    Extract text from an image using pytesseract.
    Accepts a Flask FileStorage object (e.g. from request.files).
    """
    # Read the image bytes and convert to a PIL image
    image_bytes = file_storage.read()
    image = Image.open(io.BytesIO(image_bytes))

    # Use Tesseract to do OCR on the image
    text = pytesseract.image_to_string(image, lang='tha+eng')

    return text.strip()

def extract_txt_file(file):
    """Try to read file as UTF-8 text"""
    try:
        return file.read().decode('utf-8', errors='ignore')
    except Exception as e:
        return ""

def extract_docx_text(file):
    doc = docx.Document(file)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_pptx_text(file):
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return '\n'.join(texts)

def extract_excel_text(file):
    """Extract text from .xlsx Excel file"""
    wb = openpyxl.load_workbook(file, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell) for cell in row if cell is not None]
            if row_text:
                texts.append(" ".join(row_text))
    return "\n".join(texts)

def extract_xls_text(file):
    """Extract text from .xls Excel file"""
    book = xlrd.open_workbook(file_contents=file.read())
    texts = []
    for sheet in book.sheets():
        for row_idx in range(sheet.nrows):
            row = sheet.row(row_idx)
            row_text = [str(cell.value) for cell in row if cell.value != ""]
            if row_text:
                texts.append(" ".join(row_text))
    return "\n".join(texts)


def encode_text_for_embedding(text: str) -> list[float]:
    """
    Convert text into an embedding vector using SentenceTransformer.
    """
    embedding = model.encode(text)
    return embedding.tolist()


def clean_text(input_text: str) -> str:
    """
    Remove NUL (0x00) and other control characters from text.
    """
    if input_text is None:
        return ""
    # Remove NUL characters
    cleaned = input_text.replace("\x00", "")
    # Optionally, remove all non-printable control chars except common whitespace (\n, \r, \t)
    cleaned = re.sub(r"[^\x20-\x7E\n\r\t]", "", cleaned)
    return cleaned

def save_vector_to_db(user_id, chat_history_id, file_name, text, embedding):
    """
    Save embedding to the document_embeddings table in PostgreSQL using pgvector.
    """
    # Convert Python list to PostgreSQL vector literal (e.g., '[0.1, 0.2, 0.3]')
    vector_literal = f"[{', '.join(map(str, embedding))}]"

    try:
        # Clean text before inserting
        text = clean_text(text)

        conn = psycopg2.connect(
            dbname=os.getenv("PGDATABASE", "ai_agent"),
            user=os.getenv("PGUSER", "athip"),
            password=os.getenv("PGPASSWORD", "123456"),
            host=os.getenv("PGHOST", "localhost"),
            port=os.getenv("PGPORT", "5432"),
        )
        cur = conn.cursor()

        query = """
            INSERT INTO document_embeddings (user_id, chat_history_id, file_name, extracted_text, embedding)
            VALUES (%s, %s, %s, %s, %s)
        """

        cur.execute(query, (user_id, chat_history_id, file_name, text, vector_literal))

        conn.commit()
        cur.close()
        conn.close()
        print("✅ Vector saved to database.")
    except Exception as e:
        print("❌ Failed to save vector:", e)


def search_similar_documents_by_chat(query_text: str, user_id: int, chat_history_id: int, top_k: int = 5):
    """
    Search for the top-k most similar document embeddings for a given chat_history_id and user_id.
    """
    # Step 1: Encode the query text to a vector
    query_embedding = model.encode(query_text).tolist()
    query_vector = f"[{', '.join(map(str, query_embedding))}]"

    try:
        conn = psycopg2.connect(
            dbname=os.getenv("PGDATABASE", "ai_agent"),
            user=os.getenv("PGUSER", "athip"),
            password=os.getenv("PGPASSWORD", "123456"),
            host=os.getenv("PGHOST", "localhost"),
            port=os.getenv("PGPORT", "5432"),
        )
        cur = conn.cursor()

        # Step 2: Search within same user and same chat
        query = """
            SELECT 
                id, file_name, extracted_text, embedding <=> %s AS distance
            FROM document_embeddings
            WHERE user_id = %s AND chat_history_id = %s
            ORDER BY embedding <=> %s
            LIMIT %s
        """

        cur.execute(query, (query_vector, user_id, chat_history_id, query_vector, top_k))
        results = cur.fetchall()

        cur.close()
        conn.close()

        # Step 3: Return results
        return [
            {
                'id': row[0],
                'file_name': row[1],
                'text': row[2],
                'distance': row[3]
            }
            for row in results
        ]

    except Exception as e:
        print("❌ Failed to perform similarity search:", e)
        return []
    
import os

class EditedFileSystem:
    def __init__(self):
        # Define a base directory for file operations
        # You might want to make this configurable
        self.base_dir = "managed_files"
        if not os.path.exists(self.base_dir):
            os.makedirs(self.base_dir)

    def _get_full_path(self, file_name: str) -> str:
        """Helper to get the full path of a file within the managed directory."""
        return os.path.join(self.base_dir, file_name)

    # --- Read file operations ---
    def read_line(self, file_name: str, start_line: int, end_line: int) -> tuple[list, str]:
        """
        Reads specific lines from a file, prepending line numbers to each.

        Args:
            file_name (str): The name of the file.
            start_line (int): The starting line number (1-indexed).
            end_line (int): The ending line number (1-indexed).

        Returns:
            tuple[list, str]: A tuple containing a list of lines (with line numbers)
                              and an error message (or empty string).
        """
        full_path = self._get_full_path(file_name)
        lines = []
        try:
            with open(full_path, 'r') as f:
                for i, line in enumerate(f, 1):
                    if start_line <= i <= end_line:
                        # clean_line = line.rstrip('\n')
                        lines.append(f"{i}: {line.rstrip()}")
                    elif i > end_line:
                        break # Optimization: stop reading once past the end_line
            return lines, ""
        except FileNotFoundError:
            return [], f"File '{file_name}' not found."
        except Exception as e:
            return [], f"An error occurred while reading file: {e}"

    def read_all(self, file_name: str) -> tuple[list, str]:
        """
        Reads all lines from a file, prepending line numbers to each.

        Args:
            file_name (str): The name of the file.

        Returns:
            tuple[list, str]: A tuple containing a list of lines (with line numbers)
                              and an error message (or empty string).
        """
        full_path = self._get_full_path(file_name)
        lines = []
        try:
            with open(full_path, 'r') as f:
                # Enumerate to get line numbers and prepend them
                lines = [f"{i}: {line.rstrip()}" for i, line in enumerate(f, 1)]
            return lines, ""
        except FileNotFoundError:
            return [], f"File '{file_name}' not found."
        except Exception as e:
            return [], f"An error occurred while reading file: {e}"

    def read_start_until_line_n(self, file_name: str, end_line: int) -> tuple[list, str]:
        """
        Reads from the start of the file up to a specified line number,
        prepending line numbers to each.

        Args:
            file_name (str): The name of the file.
            end_line (int): The ending line number (1-indexed).

        Returns:
            tuple[list, str]: A tuple containing a list of lines (with line numbers)
                              and an error message (or empty string).
        """
        # This method already calls read_line, which now handles line numbering
        return self.read_line(file_name, 1, end_line)

    # --- Edit file operations ---
    def edit_line(self, file_name: str, text: str, start_line: int, end_line: int) -> str:
        """
        Edits specific lines in a file by replacing them with the provided text.

        Args:
            file_name (str): The name of the file.
            text (str): The new text to insert.
            start_line (int): The starting line number (1-indexed) to replace.
            end_line (int): The ending line number (1-indexed) to replace.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = self._get_full_path(file_name)
        try:
            with open(full_path, 'r') as f:
                lines = f.readlines()

            start_idx = start_line - 1
            end_idx = end_line - 1

            if not (0 <= start_idx <= end_idx < len(lines)):
                return f"Line range {start_line}-{end_line} is out of bounds for file '{file_name}'."

            new_lines_to_insert = [line + '\n' for line in text.splitlines()]
            if not new_lines_to_insert:
                new_lines_to_insert = ['\n'] # Ensure at least one newline for empty text

            new_file_content = lines[:start_idx] + new_lines_to_insert + lines[end_idx + 1:]

            with open(full_path, 'w') as f:
                f.writelines(new_file_content)
            return ""
        except FileNotFoundError:
            return f"File '{file_name}' not found for editing."
        except Exception as e:
            return f"An error occurred while editing file: {e}"

    def edit_all(self, file_name: str, text: str) -> str:
        """
        Replaces the entire content of a file with the provided text.

        Args:
            file_name (str): The name of the file.
            text (str): The new text to write to the file.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = self._get_full_path(file_name)
        try:
            with open(full_path, 'w') as f:
                f.write(text)
            return ""
        except Exception as e:
            return f"An error occurred while editing all of file: {e}"

    # --- Create new file operations ---
    def create_new_file_only(self, file_name: str) -> str:
        """
        Creates a new empty file.

        Args:
            file_name (str): The name of the file to create.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = self._get_full_path(file_name)
        if os.path.exists(full_path):
            return f"File '{file_name}' already exists."
        try:
            with open(full_path, 'x') as f:
                pass
            return ""
        except FileExistsError: # Redundant due to above check, but good for robustness
            return f"File '{file_name}' already exists unexpectedly."
        except Exception as e:
            return f"An error occurred while creating new file: {e}"

    def create_new_file_and_text(self, file_name: str, text: str) -> str:
        """
        Creates a new file and writes the provided text to it.
        Ensures the directory path exists before creating the file.

        Args:
            file_name (str): The name of the file to create (can include subdirectories).
            text (str): The text to write into the new file.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = self._get_full_path(file_name)
        
        # Extract the directory path from the full file path
        directory = os.path.dirname(full_path)
        
        try:
            # Create the directory if it doesn't exist (and any parent directories)
            # exist_ok=True prevents an error if the directory already exists
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True) # This is the key line for cascade creation

            with open(full_path, 'w') as f:
                f.write(text)
            return ""
        except Exception as e:
            return f"An error occurred while creating new file with text: {e}"

    def delete_file(self, file_name: str) -> str:
        """
        Deletes a file from the managed directory.

        Args:
            file_name (str): The name of the file to delete.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = self._get_full_path(file_name)
        try:
            if os.path.exists(full_path):
                os.remove(full_path)
                return ""
            else:
                return f"File '{file_name}' not found."
        except Exception as e:
            return f"An error occurred while deleting file: {e}"

    def list_files(self) -> tuple[list, str]:
        """
        Lists all files in the managed directory.

        Returns:
            tuple[list, str]: A tuple containing a list of file names and an error message (or empty string).
        """
        try:
            files = [f for f in os.listdir(self.base_dir) if os.path.isfile(self._get_full_path(f))]
            return files, ""
        except Exception as e:
            return [], f"An error occurred while listing files: {e}"

    # --- Folder operations ---
    def create_folder(self, folder_name: str) -> str:
        """
        Creates a new folder within the managed directory.

        Args:
            folder_name (str): The name of the folder to create.

        Returns:
            str: An empty string if successful, otherwise an error message.
        """
        full_path = os.path.join(self.base_dir, folder_name)
        try:
            if os.path.exists(full_path):
                return f"Folder '{folder_name}' already exists."
            os.makedirs(full_path)
            return ""
        except Exception as e:
            return f"An error occurred while creating folder: {e}"